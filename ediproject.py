import pyttsx3
import speech_recognition as sr
import datetime
import os
import pyaudio
from pptx import Presentation
from docx import Document
from pptx.util import Inches, Pt
 
engine = pyttsx3.init('sapi5')
voices = engine.getProperty('voices') #print(voices[0].id)
engine.setProperty('voice', voices[1].id)


def speak(audio):
    engine.say(audio)
    engine.runAndWait()

def wishMe():
    hour = int(datetime.datetime.now().hour)
    if hour>=0 and hour<12:
        speak("Good Morning!")

    elif hour>=12 and hour<18:
        speak("Good Afternoon!")

    else :
        speak("Good Evening!")

    speak("I am Alice, tell me what u want, I am there to help you!")

def takeCommand():
    r = sr.Recognizer()
    with sr.Microphone() as source:
        print("Listening...")
        r.pause_threshold = 1
        audio = r.listen(source)

    try:
        print("recognizing...")
        query = r.recognize_google(audio, language='en-in')
        print(f"user said: {query}\n")
        

    except Exception as e:
        #print(e)
        print("please say again...")
        return "None"
    return query    

if __name__ == "__main__":
    wishMe()
    file_name = ""
    while True:
   
        query = takeCommand().lower()
        if 'how are you' in query:
            speak("I am fine, Thank you")
            speak("How are you, Sir")
        
        elif 'fine' in query or "good" in query:
            speak("It's good to know that your fine")

        elif 'the time' in query:
            strTime = datetime.datetime.now().strftime("%I:%M %p")
            speak(f"Sir, the time is {strTime}")
        
        elif 'change name' in query:
            speak("what would u like to call me sir? ")
            askname = takeCommand()
            speak("Amazing!, I liked your name!")


        elif 'default location' in query:
            speak("opening default location...")
            codePath ="D:\\SY DOC"
            os.startfile(codePath)    
        
        elif 'open vs code' in query:
            speak("opening vscode...")
            codePath = "C:\\Users\\Lenovo\\AppData\\Local\\Programs\\Microsoft VS Code\\Code.exe"
            os.startfile(codePath)

        elif 'open wordpad' in query:
            speak("opening wordpad....")
            codePath = "C:\\ProgramData\\Microsoft\\Windows\\Start Menu\\Programs\\Accessories\\Wordpad"
            os.startfile(codePath)

        
        elif 'open powerpoint' in query:
            speak("opening powerpoint....")
            codePath = "C:\\Program Files (x86)\\Microsoft Office\\root\\Office16\\POWERPNT.EXE"
            os.startfile(codePath)
        
        elif 'open excel' in query:
            speak("opening excel....")
            codePath = "C:\\Program Files (x86)\\Microsoft Office\\root\\Office16\\EXCEL.EXE"
            os.startfile(codePath)

        elif 'open notepad' in query:
            speak("opening notepad....")
            codePath = "C:\\ProgramData\\Microsoft\\Windows\\Start Menu\\Programs\\Accessories\\Notepad"
            os.startfile(codePath)

        
        elif 'open paint' in query:
            speak("opening paint....")
            codePath = "C:\\ProgramData\\Microsoft\\Windows\\Start Menu\\Programs\\Accessories\\Paint"
            os.startfile(codePath)

        elif 'open microsoft edge' in query:
            speak("opening Microsoft Edge....")
            codePath = "C:\\Program Files (x86)\\Microsoft\\Edge\\Application\\msedge.exe"
            os.startfile(codePath)

        elif 'open chrome' in query:
            speak("opening chrome....")
            codePath = "C:\\Program Files (x86)\\Google\\Chrome\\Application\\chrome.exe"
            os.startfile(codePath)

        elif 'open command prompt' in query:
            speak("opening command prompt....")
            codePath = "C:\\Users\\Lenovo\\AppData\\Roaming\\Microsoft\\Windows\\Start Menu\\Programs\\System Tools\\Command Prompt"
            os.startfile(codePath)
        

        elif 'play music' in query or "play song" in query:
            speak("Here you go with music")
            # music_dir = "G:\\Song"
            music_dir = "C:\\Users\\Lenovo\\Music"
            songs = os.listdir(music_dir)
            print(songs)
            speak("tell me the number of your song")
            num = int(takeCommand())
            random = os.startfile(os.path.join(music_dir, songs[num]))

        elif "write a note" in query:
            speak("What should i write, sir")
            note = takeCommand()
            file = open('jarvis.txt', 'w')
            speak("Sir, Should i include date and time")
            snfm = takeCommand()
            if 'yes' in snfm or 'sure' in snfm:
                strTime = datetime.datetime.now().strftime("%I:%M %p")
                file.write(strTime)
                file.write(" :- ")
                file.write(note)
            else:
                file.write(note)
         
        elif "show note" in query:
            speak("Showing Notes")
            file = open("jarvis.txt", "r")
            print(file.read())
            speak(file.read(6))

        elif "create powerpoint presentation" in query:
            root = Presentation()   # presentation obj
            first_layer = root.slide_layouts[0] # 0 will contian title and subtitle layout 
            slide = root.slides.add_slide(first_layer) #creating slide
            speak("Sir, what title should I add?")
            title = takeCommand().upper()
            slide.shapes.title.text = title
          
            speak("Sir, what subtitle should I add?")
            subtitle = takeCommand().capitalize()
            slide.placeholders[1].text = subtitle
          
            speak("sir, you want to add more slides?")
            command = takeCommand()
            while "no" not in command or "yes" in command:
                second_slide_layout = root.slide_layouts[6]  #6 for blank slide
                slide = root.slides.add_slide(second_slide_layout)
                speak("New slide added")
                speak("What would you like to add in it?")
                print("Text\nImage\nTable\n")
                command = takeCommand()
                if "text" in command:
                    left = top = width = height = Inches(1) 
                    txBox = slide.shapes.add_textbox(left, top,width, height) #creating the textbox
                    tf = txBox.text_frame   #creating the text frame
                    command = "yes"
                    while "no" not in command:
                        speak("What Text you want to add")
                        p = tf.add_paragraph()
                        p.text = takeCommand().capitalize()
                        speak("Would you like to change font size")
                        command = takeCommand().lower()
                        if "yes" in command:
                            speak("Tell me the new font size")
                            command = int(takeCommand())
                            p.font.size = Pt(command)
                        else:
                            p.font.size = Pt(20)

                        speak("Would you like to bold your font or would you like to change your font style?")
                        command = takeCommand()
                        if "yes" in command:
                            speak("font or style??")
                            command = takeCommand()
                            p.font.Calibri = True
                            if "style" in command:
                                print("italic\nVerdana\nCalibri\nPalatino\nTahoma\nGeorgia\nGill Sans\nCorbel\nSegoe\n")
                                speak("Choose the style...")
                                command = takeCommand()
                                p.font.command = True
                        
                            elif "font" in command:
                               p.font.bold = True
                            
                        speak("Want to add more text in it?")
                        command = takeCommand().lower()
                
                elif "image" in command:
                   # slide = root.slides.add_slide(second_slide_layout)
                    left = top = Inches(1)
                    height = Inches(5)
                    speak("Name of the image")
                    img_path = takeCommand()
                    pic = slide.shapes.add_picture(img_path+".jpg",left, top,height=height)
                    speak("Image added")
                
                elif "table" in command:
                    x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
                    shape = slide.shapes.add_table(3, 4, x, y, cx, cy)
                    speak("Table added")
                
                speak("Would you want to add more slides?")
                command = takeCommand()



            
            speak("Tell me the name for this file?")
            file_name = takeCommand()
            root.save(file_name+".pptx")
            speak("File is saved")
           
        elif "open my presentation file" in query:
            speak("What is the name of the file")
            file_name = takeCommand()
            speak("opening"+file_name+"....")
            codePath = file_name+".pptx"
            os.startfile(codePath)
        
        elif 'create a word file' in query:
            doc = Document()     #creating obj
            speak("creating a word file...")
            speak("sir, what heading should I add?...")
            heading = takeCommand().upper()
            doc.add_heading(heading,level=0)
            speak("What should I add in The Paragraph...")
            para = takeCommand().capitalize()
            doc.add_paragraph(para,style = "List Bullet")
            speak("Text added")
            speak("Would you like to  continue...")
            command = takeCommand().lower()
            while "no"  not in command or "yes" in command :
                speak("What should I add in the next line")
                print("\nparagraph\nImage\n")
                command = takeCommand().lower()
                if "paragraph" in command:
                    speak("Please add a paragraph!!")
                    para2 = takeCommand()
                    doc.add_paragraph(para2,style = "List Bullet")
                    speak("text added ")

                elif "image" in command:
                    speak("Please tell the name of your image!")
                    img = takeCommand()
                    doc.add_picture(img+".jpg")
                    speak("Image added ")
                
                speak("Would you like to  continue...")
                command = takeCommand()
            
            speak("Please tell me the name for this file?")
            file_name = takeCommand()
            doc.save(file_name+".docx")
            speak("File is saved")

        elif "open word file" in query:
            speak("What the name of the file?")
            file_name = takeCommand()
            speak("Opening"+file_name+"...")
            codePath = file_name+".docx"
            os.startfile(codePath)

        elif "create a folder" in query:
            speak("Tell me the name of the folder:")
            directory = takeCommand()
            parent_dir = "D:/SY DOC/"
            path = os.path.join(parent_dir,directory)
            os.mkdir(path)
            print("Directory '% s' created" % directory)
            speak("do you want to open the same")
            command = takeCommand()
            if "yes" in command:
                os.startfile(path)

        elif "delete a folder" in query:
            speak("Tell me the name of the folder:")
            directory = takeCommand()
            parent_dir = "D:/SY DOC/"
            path = os.path.join(parent_dir,directory)
            os.rmdir(path) 
            print("Directory '% s' deleted" % directory)

        elif "open a folder" in query:
            speak("Tell me the name of the folder:")
            directory = takeCommand()
            parent_dir = "D:/SY DOC/"
            path = os.path.join(parent_dir,directory)
            os.startfile(path)


        elif 'exit' in query:
            speak("I am very thankful to share few things with you! Have a great day, sir! ")
            exit() 
        
        else:
            speak("Please say again")
    

 

        


